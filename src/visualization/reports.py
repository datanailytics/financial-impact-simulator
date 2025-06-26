"""
Reports Module - Professional Financial Report Generation
Author: Naiara Rodríguez Solano
Email: datanailytics@outlook.com
GitHub: https://github.com/datanailytics
Portfolio: https://datanailytics.github.io

This module provides comprehensive report generation capabilities for financial
analysis, supporting multiple formats including PDF, Excel, PowerPoint, and HTML.
"""

import pandas as pd
import numpy as np
from typing import Dict, List, Tuple, Optional, Union, Any
from datetime import datetime, timedelta
import os
from pathlib import Path
import json
import logging
from dataclasses import dataclass, field
from enum import Enum

# Report generation libraries
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_CENTER, TA_RIGHT, TA_LEFT
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.linecharts import HorizontalLineChart
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.graphics.charts.piecharts import Pie

import xlsxwriter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import seaborn as sns
from jinja2 import Environment, FileSystemLoader, Template

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class ReportFormat(Enum):
    """Available report formats."""
    PDF = "pdf"
    EXCEL = "excel"
    HTML = "html"
    POWERPOINT = "powerpoint"
    WORD = "word"
    MARKDOWN = "markdown"
    JSON = "json"


class ReportType(Enum):
    """Types of financial reports."""
    PORTFOLIO_SUMMARY = "portfolio_summary"
    RISK_ANALYSIS = "risk_analysis"
    PERFORMANCE = "performance"
    COMPLIANCE = "compliance"
    TRADING = "trading"
    RESEARCH = "research"
    CUSTOM = "custom"


@dataclass
class ReportConfig:
    """Configuration for report generation."""
    title: str
    subtitle: Optional[str] = None
    author: str = "Naiara Rodríguez Solano"
    company: str = "DataNailytics"
    report_type: ReportType = ReportType.PORTFOLIO_SUMMARY
    format: ReportFormat = ReportFormat.PDF
    include_charts: bool = True
    include_tables: bool = True
    include_summary: bool = True
    include_disclaimer: bool = True
    logo_path: Optional[str] = None
    template_path: Optional[str] = None
    output_path: str = "./reports"
    
    def __post_init__(self):
        """Create output directory if it doesn't exist."""
        Path(self.output_path).mkdir(parents=True, exist_ok=True)


@dataclass
class ReportSection:
    """Represents a section in the report."""
    title: str
    content: Union[str, pd.DataFrame, Dict[str, Any]]
    section_type: str = "text"  # text, table, chart, image
    metadata: Dict[str, Any] = field(default_factory=dict)


class BaseReportGenerator:
    """Base class for report generators."""
    
    def __init__(self, config: ReportConfig):
        """
        Initialize report generator.
        
        Args:
            config: Report configuration
        """
        self.config = config
        self.sections: List[ReportSection] = []
        self.metadata = {
            "created_at": datetime.now(),
            "author": config.author,
            "company": config.company,
            "report_type": config.report_type.value
        }
    
    def add_section(self, section: ReportSection):
        """Add a section to the report."""
        self.sections.append(section)
    
    def add_title_page(self):
        """Add title page to report."""
        title_content = {
            "title": self.config.title,
            "subtitle": self.config.subtitle,
            "author": self.config.author,
            "company": self.config.company,
            "date": datetime.now().strftime("%B %d, %Y")
        }
        
        self.add_section(ReportSection(
            title="Title Page",
            content=title_content,
            section_type="title"
        ))
    
    def add_executive_summary(self, summary_data: Dict[str, Any]):
        """Add executive summary section."""
        self.add_section(ReportSection(
            title="Executive Summary",
            content=summary_data,
            section_type="summary"
        ))
    
    def add_disclaimer(self):
        """Add standard disclaimer."""
        disclaimer_text = """
        DISCLAIMER: This report is provided for informational purposes only and does not constitute 
        investment advice, financial advice, trading advice, or any other sort of advice. You should 
        not treat any of the report's content as such. The author does not recommend that any 
        cryptocurrency or financial instrument should be bought, sold, or held by you. Do conduct 
        your own due diligence and consult your financial advisor before making any investment decisions.
        
        Past performance is not indicative of future results. All investments carry risk, including 
        the potential loss of principal.
        """
        
        self.add_section(ReportSection(
            title="Disclaimer",
            content=disclaimer_text,
            section_type="disclaimer"
        ))
    
    def generate(self) -> str:
        """Generate the report."""
        raise NotImplementedError("Subclasses must implement generate method")


class PDFReportGenerator(BaseReportGenerator):
    """Generate PDF reports using ReportLab."""
    
    def __init__(self, config: ReportConfig):
        """Initialize PDF report generator."""
        super().__init__(config)
        self.styles = getSampleStyleSheet()
        self._setup_custom_styles()
    
    def _setup_custom_styles(self):
        """Setup custom paragraph styles."""
        # Title style
        self.styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=self.styles['Title'],
            fontSize=24,
            textColor=colors.HexColor('#1a1a2e'),
            spaceAfter=30,
            alignment=TA_CENTER
        ))
        
        # Subtitle style
        self.styles.add(ParagraphStyle(
            name='CustomSubtitle',
            parent=self.styles['Title'],
            fontSize=18,
            textColor=colors.HexColor('#16213e'),
            spaceAfter=20,
            alignment=TA_CENTER
        ))
        
        # Section header style
        self.styles.add(ParagraphStyle(
            name='SectionHeader',
            parent=self.styles['Heading1'],
            fontSize=16,
            textColor=colors.HexColor('#0f3460'),
            spaceAfter=12,
            spaceBefore=20
        ))
    
    def generate(self) -> str:
        """Generate PDF report."""
        filename = f"{self.config.title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        filepath = os.path.join(self.config.output_path, filename)
        
        doc = SimpleDocTemplate(
            filepath,
            pagesize=letter,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=18
        )
        
        story = []
        
        # Add sections
        for section in self.sections:
            if section.section_type == "title":
                story.extend(self._create_title_page(section.content))
            elif section.section_type == "text":
                story.extend(self._create_text_section(section))
            elif section.section_type == "table":
                story.extend(self._create_table_section(section))
            elif section.section_type == "chart":
                story.extend(self._create_chart_section(section))
            elif section.section_type == "summary":
                story.extend(self._create_summary_section(section))
            
            story.append(Spacer(1, 0.2*inch))
        
        # Build PDF
        doc.build(story)
        logger.info(f"PDF report generated: {filepath}")
        
        return filepath
    
    def _create_title_page(self, content: Dict[str, Any]) -> List:
        """Create title page elements."""
        elements = []
        
        # Add logo if provided
        if self.config.logo_path and os.path.exists(self.config.logo_path):
            logo = Image(self.config.logo_path, width=2*inch, height=1*inch)
            elements.append(logo)
            elements.append(Spacer(1, 0.5*inch))
        
        # Title
        elements.append(Paragraph(content['title'], self.styles['CustomTitle']))
        
        # Subtitle
        if content.get('subtitle'):
            elements.append(Paragraph(content['subtitle'], self.styles['CustomSubtitle']))
        
        elements.append(Spacer(1, 1*inch))
        
        # Author and company
        elements.append(Paragraph(content['author'], self.styles['Normal']))
        elements.append(Paragraph(content['company'], self.styles['Normal']))
        elements.append(Spacer(1, 0.5*inch))
        
        # Date
        elements.append(Paragraph(content['date'], self.styles['Normal']))
        
        elements.append(PageBreak())
        
        return elements
    
    def _create_text_section(self, section: ReportSection) -> List:
        """Create text section elements."""
        elements = []
        
        # Section header
        elements.append(Paragraph(section.title, self.styles['SectionHeader']))
        
        # Content
        if isinstance(section.content, str):
            # Split by paragraphs
            for para in section.content.split('\n\n'):
                if para.strip():
                    elements.append(Paragraph(para.strip(), self.styles['Normal']))
                    elements.append(Spacer(1, 0.1*inch))
        
        return elements
    
    def _create_table_section(self, section: ReportSection) -> List:
        """Create table section elements."""
        elements = []
        
        # Section header
        elements.append(Paragraph(section.title, self.styles['SectionHeader']))
        
        if isinstance(section.content, pd.DataFrame):
            df = section.content
            
            # Create table data
            data = [df.columns.tolist()]  # Header
            data.extend(df.values.tolist())  # Data rows
            
            # Create table
            table = Table(data)
            
            # Apply table style
            style = TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1a1a2e')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                ('FONTSIZE', (0, 1), (-1, -1), 10),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey])
            ])
            
            table.setStyle(style)
            elements.append(table)
        
        return elements
    
    def _create_chart_section(self, section: ReportSection) -> List:
        """Create chart section elements."""
        elements = []
        
        # Section header
        elements.append(Paragraph(section.title, self.styles['SectionHeader']))
        
        # Create and add chart
        if section.metadata.get('chart_type') == 'line':
            chart = self._create_line_chart(section.content)
        elif section.metadata.get('chart_type') == 'bar':
            chart = self._create_bar_chart(section.content)
        elif section.metadata.get('chart_type') == 'pie':
            chart = self._create_pie_chart(section.content)
        else:
            # Default to creating an image from matplotlib
            chart_path = self._create_matplotlib_chart(section.content, section.metadata)
            chart = Image(chart_path, width=6*inch, height=4*inch)
        
        elements.append(chart)
        
        return elements
    
    def _create_summary_section(self, section: ReportSection) -> List:
        """Create executive summary section."""
        elements = []
        
        # Section header
        elements.append(Paragraph(section.title, self.styles['SectionHeader']))
        
        if isinstance(section.content, dict):
            for key, value in section.content.items():
                # Create key-value pairs
                text = f"<b>{key}:</b> {value}"
                elements.append(Paragraph(text, self.styles['Normal']))
                elements.append(Spacer(1, 0.05*inch))
        
        return elements
    
    def _create_matplotlib_chart(self, data: Any, metadata: Dict[str, Any]) -> str:
        """Create a chart using matplotlib and save as image."""
        plt.figure(figsize=(10, 6))
        
        chart_type = metadata.get('chart_type', 'line')
        
        if isinstance(data, pd.DataFrame):
            if chart_type == 'line':
                data.plot(kind='line')
            elif chart_type == 'bar':
                data.plot(kind='bar')
            elif chart_type == 'pie' and len(data.columns) == 1:
                data.plot(kind='pie', y=data.columns[0])
        
        plt.title(metadata.get('title', ''))
        plt.xlabel(metadata.get('xlabel', ''))
        plt.ylabel(metadata.get('ylabel', ''))
        
        # Save chart
        chart_path = os.path.join(self.config.output_path, f"chart_{datetime.now().timestamp()}.png")
        plt.savefig(chart_path, dpi=300, bbox_inches='tight')
        plt.close()
        
        return chart_path


class ExcelReportGenerator(BaseReportGenerator):
    """Generate Excel reports using XlsxWriter."""
    
    def generate(self) -> str:
        """Generate Excel report."""
        filename = f"{self.config.title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        filepath = os.path.join(self.config.output_path, filename)
        
        # Create workbook
        workbook = xlsxwriter.Workbook(filepath)
        
        # Add formats
        formats = self._create_formats(workbook)
        
        # Add summary sheet
        summary_sheet = workbook.add_worksheet('Summary')
        self._write_summary_sheet(summary_sheet, formats)
        
        # Add sections as separate sheets
        for section in self.sections:
            if section.section_type == "table" and isinstance(section.content, pd.DataFrame):
                sheet_name = section.title.replace(' ', '_')[:31]  # Excel sheet name limit
                worksheet = workbook.add_worksheet(sheet_name)
                self._write_dataframe_to_sheet(worksheet, section.content, formats)
            elif section.section_type == "summary" and isinstance(section.content, dict):
                sheet_name = section.title.replace(' ', '_')[:31]
                worksheet = workbook.add_worksheet(sheet_name)
                self._write_dict_to_sheet(worksheet, section.content, formats)
        
        # Add charts if enabled
        if self.config.include_charts:
            chart_sheet = workbook.add_worksheet('Charts')
            self._add_charts(workbook, chart_sheet)
        
        workbook.close()
        logger.info(f"Excel report generated: {filepath}")
        
        return filepath
    
    def _create_formats(self, workbook: xlsxwriter.Workbook) -> Dict[str, Any]:
        """Create cell formats for Excel."""
        formats = {
            'title': workbook.add_format({
                'bold': True,
                'font_size': 20,
                'font_color': '#1a1a2e',
                'align': 'center',
                'valign': 'vcenter'
            }),
            'header': workbook.add_format({
                'bold': True,
                'font_size': 12,
                'font_color': 'white',
                'bg_color': '#1a1a2e',
                'border': 1,
                'align': 'center',
                'valign': 'vcenter'
            }),
            'cell': workbook.add_format({
                'font_size': 10,
                'border': 1,
                'align': 'center',
                'valign': 'vcenter'
            }),
            'number': workbook.add_format({
                'font_size': 10,
                'border': 1,
                'align': 'right',
                'valign': 'vcenter',
                'num_format': '#,##0.00'
            }),
            'percent': workbook.add_format({
                'font_size': 10,
                'border': 1,
                'align': 'right',
                'valign': 'vcenter',
                'num_format': '0.00%'
            }),
            'date': workbook.add_format({
                'font_size': 10,
                'border': 1,
                'align': 'center',
                'valign': 'vcenter',
                'num_format': 'yyyy-mm-dd'
            })
        }
        
        return formats
    
    def _write_summary_sheet(self, worksheet, formats: Dict[str, Any]):
        """Write summary information to worksheet."""
        # Title
        worksheet.merge_range('A1:F2', self.config.title, formats['title'])
        
        # Metadata
        row = 4
        worksheet.write(row, 0, 'Report Date:', formats['header'])
        worksheet.write(row, 1, datetime.now().strftime('%Y-%m-%d'), formats['date'])
        
        row += 1
        worksheet.write(row, 0, 'Author:', formats['header'])
        worksheet.write(row, 1, self.config.author, formats['cell'])
        
        row += 1
        worksheet.write(row, 0, 'Company:', formats['header'])
        worksheet.write(row, 1, self.config.company, formats['cell'])
        
        # Adjust column widths
        worksheet.set_column('A:A', 20)
        worksheet.set_column('B:F', 15)
    
    def _write_dataframe_to_sheet(self, worksheet, df: pd.DataFrame, formats: Dict[str, Any]):
        """Write DataFrame to worksheet."""
        # Write headers
        for col_num, column in enumerate(df.columns):
            worksheet.write(0, col_num, column, formats['header'])
        
        # Write data
        for row_num, row_data in enumerate(df.values):
            for col_num, value in enumerate(row_data):
                # Determine format based on value type
                if isinstance(value, (int, float)):
                    if 'percent' in str(df.columns[col_num]).lower() or 'rate' in str(df.columns[col_num]).lower():
                        worksheet.write(row_num + 1, col_num, value, formats['percent'])
                    else:
                        worksheet.write(row_num + 1, col_num, value, formats['number'])
                elif isinstance(value, datetime):
                    worksheet.write(row_num + 1, col_num, value, formats['date'])
                else:
                    worksheet.write(row_num + 1, col_num, str(value), formats['cell'])
        
        # Autofit columns
        for col_num, column in enumerate(df.columns):
            worksheet.set_column(col_num, col_num, max(len(str(column)), 12))
    
    def _write_dict_to_sheet(self, worksheet, data: Dict[str, Any], formats: Dict[str, Any]):
        """Write dictionary data to worksheet."""
        row = 0
        for key, value in data.items():
            worksheet.write(row, 0, key, formats['header'])
            
            if isinstance(value, (int, float)):
                worksheet.write(row, 1, value, formats['number'])
            else:
                worksheet.write(row, 1, str(value), formats['cell'])
            
            row += 1
        
        # Adjust column widths
        worksheet.set_column('A:A', 30)
        worksheet.set_column('B:B', 20)
    
    def _add_charts(self, workbook: xlsxwriter.Workbook, worksheet):
        """Add charts to worksheet."""
        # Example: Add a line chart
        chart = workbook.add_chart({'type': 'line'})
        
        # Configure chart (this is a placeholder - actual data needs to be referenced)
        chart.set_title({'name': 'Performance Chart'})
        chart.set_x_axis({'name': 'Date'})
        chart.set_y_axis({'name': 'Value'})
        
        worksheet.insert_chart('A1', chart, {'x_scale': 2, 'y_scale': 2})


class PowerPointReportGenerator(BaseReportGenerator):
    """Generate PowerPoint reports using python-pptx."""
    
    def generate(self) -> str:
        """Generate PowerPoint report."""
        filename = f"{self.config.title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join(self.config.output_path, filename)
        
        # Create presentation
        prs = Presentation()
        
        # Set slide size to widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # Add title slide
        self._add_title_slide(prs)
        
        # Add content slides
        for section in self.sections:
            if section.section_type == "text":
                self._add_text_slide(prs, section)
            elif section.section_type == "table":
                self._add_table_slide(prs, section)
            elif section.section_type == "chart":
                self._add_chart_slide(prs, section)
            elif section.section_type == "summary":
                self._add_summary_slide(prs, section)
        
        # Save presentation
        prs.save(filepath)
        logger.info(f"PowerPoint report generated: {filepath}")
        
        return filepath
    
    def _add_title_slide(self, prs: Presentation):
        """Add title slide to presentation."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = self.config.title
        
        # Subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = f"{self.config.author}\n{self.config.company}\n{datetime.now().strftime('%B %d, %Y')}"
    
    def _add_text_slide(self, prs: Presentation, section: ReportSection):
        """Add text slide to presentation."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = section.title
        
        # Content
        content = slide.placeholders[1]
        if isinstance(section.content, str):
            content.text = section.content
    
    def _add_table_slide(self, prs: Presentation, section: ReportSection):
        """Add table slide to presentation."""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(12)
        height = Inches(1)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = section.title
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        
        if isinstance(section.content, pd.DataFrame):
            df = section.content
            
            # Add table
            rows, cols = df.shape
            rows += 1  # Add header row
            
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(12)
            height = Inches(5)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Set column headers
            for col_idx, col_name in enumerate(df.columns):
                table.cell(0, col_idx).text = str(col_name)
                table.cell(0, col_idx).fill.solid()
                table.cell(0, col_idx).fill.fore_color.rgb = RGBColor(26, 26, 46)
            
            # Add data
            for row_idx, row_data in enumerate(df.values):
                for col_idx, value in enumerate(row_data):
                    table.cell(row_idx + 1, col_idx).text = str(value)
    
    def _add_chart_slide(self, prs: Presentation, section: ReportSection):
        """Add chart slide to presentation."""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(12)
        height = Inches(1)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = section.title
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        
        # Create and add chart image
        chart_path = self._create_matplotlib_chart(section.content, section.metadata)
        
        left = Inches(1)
        top = Inches(2)
        slide.shapes.add_picture(chart_path, left, top, height=Inches(5))
    
    def _add_summary_slide(self, prs: Presentation, section: ReportSection):
        """Add summary slide to presentation."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = section.title
        
        # Content
        content = slide.placeholders[1]
        text_frame = content.text_frame
        
        if isinstance(section.content, dict):
            for key, value in section.content.items():
                p = text_frame.add_paragraph()
                p.text = f"• {key}: {value}"
                p.level = 0
    
    def _create_matplotlib_chart(self, data: Any, metadata: Dict[str, Any]) -> str:
        """Create a chart using matplotlib and save as image."""
        plt.figure(figsize=(10, 6))
        
        chart_type = metadata.get('chart_type', 'line')
        
        if isinstance(data, pd.DataFrame):
            if chart_type == 'line':
                data.plot(kind='line')
            elif chart_type == 'bar':
                data.plot(kind='bar')
            elif chart_type == 'pie' and len(data.columns) == 1:
                data.plot(kind='pie', y=data.columns[0])
        
        plt.title(metadata.get('title', ''))
        plt.xlabel(metadata.get('xlabel', ''))
        plt.ylabel(metadata.get('ylabel', ''))
        
        # Save chart
        chart_path = os.path.join(self.config.output_path, f"chart_{datetime.now().timestamp()}.png")
        plt.savefig(chart_path, dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()
        
        return chart_path


class HTMLReportGenerator(BaseReportGenerator):
    """Generate HTML reports using Jinja2 templates."""
    
    def __init__(self, config: ReportConfig):
        """Initialize HTML report generator."""
        super().__init__(config)
        
        # Setup Jinja2 environment
        if config.template_path and os.path.exists(config.template_path):
            self.env = Environment(loader=FileSystemLoader(config.template_path))
        else:
            # Use default template
            self.env = Environment()
            self.env.from_string(self._get_default_template())
    
    def generate(self) -> str:
        """Generate HTML report."""
        filename = f"{self.config.title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
        filepath = os.path.join(self.config.output_path, filename)
        
        # Prepare template data
        template_data = {
            'title': self.config.title,
            'subtitle': self.config.subtitle,
            'author': self.config.author,
            'company': self.config.company,
            'date': datetime.now().strftime('%B %d, %Y'),
            'sections': self._prepare_sections_for_template(),
            'metadata': self.metadata
        }
        
        # Render template
        if self.config.template_path:
            template = self.env.get_template('report_template.html')
        else:
            template = self.env.from_string(self._get_default_template())
        
        html_content = template.render(**template_data)
        
        # Write to file
        with open(filepath, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        logger.info(f"HTML report generated: {filepath}")
        return filepath
    
    def _prepare_sections_for_template(self) -> List[Dict[str, Any]]:
        """Prepare sections for template rendering."""
        prepared_sections = []
        
        for section in self.sections:
            prepared = {
                'title': section.title,
                'type': section.section_type,
                'content': None
            }
            
            if section.section_type == "text":
                prepared['content'] = section.content
            elif section.section_type == "table" and isinstance(section.content, pd.DataFrame):
                prepared['content'] = section.content.to_html(classes='table table-striped', index=False)
            elif section.section_type == "summary" and isinstance(section.content, dict):
                prepared['content'] = section.content
            elif section.section_type == "chart":
                # Convert chart to base64 image
                chart_path = self._create_matplotlib_chart(section.content, section.metadata)
                with open(chart_path, 'rb') as f:
                    import base64
                    chart_data = base64.b64encode(f.read()).decode()
                prepared['content'] = f'data:image/png;base64,{chart_data}'
            
            prepared_sections.append(prepared)
        
        return prepared_sections
    
    def _get_default_template(self) -> str:
        """Get default HTML template."""
        return '''
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{{ title }}</title>
    <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.1.3/dist/css/bootstrap.min.css" rel="stylesheet">
    <style>
        body {
            font-family: 'Arial', sans-serif;
            line-height: 1.6;
            color: #333;
        }
        .header {
            background-color: #1a1a2e;
            color: white;
            padding: 2rem 0;
            margin-bottom: 2rem;
        }
        .section {
            margin-bottom: 3rem;
        }
        .table {
            margin-top: 1rem;
        }
        .chart-container {
            text-align: center;
            margin: 2rem 0;
        }
        .chart-container img {
            max-width: 100%;
            height: auto;
        }
        .footer {
            background-color: #f8f9fa;
            padding: 2rem 0;
            margin-top: 3rem;
            text-align: center;
        }
        .summary-item {
            padding: 0.5rem 0;
            border-bottom: 1px solid #eee;
        }
        @media print {
            .page-break {
                page-break-after: always;
            }
        }
    </style>
</head>
<body>
    <div class="header">
        <div class="container">
            <h1 class="text-center">{{ title }}</h1>
            {% if subtitle %}
            <h3 class="text-center">{{ subtitle }}</h3>
            {% endif %}
            <p class="text-center mb-0">{{ author }} | {{ company }}</p>
            <p class="text-center">{{ date }}</p>
        </div>
    </div>
    
    <div class="container">
        {% for section in sections %}
        <div class="section">
            <h2>{{ section.title }}</h2>
            
            {% if section.type == "text" %}
                <div class="content">{{ section.content | safe }}</div>
            
            {% elif section.type == "table" %}
                {{ section.content | safe }}
            
            {% elif section.type == "summary" %}
                <div class="summary">
                    {% for key, value in section.content.items() %}
                    <div class="summary-item">
                        <strong>{{ key }}:</strong> {{ value }}
                    </div>
                    {% endfor %}
                </div>
            
            {% elif section.type == "chart" %}
                <div class="chart-container">
                    <img src="{{ section.content }}" alt="{{ section.title }}">
                </div>
            {% endif %}
        </div>
        {% endfor %}
    </div>
    
    <div class="footer">
        <div class="container">
            <p class="mb-0">Generated by DataNailytics Financial Analytics Platform</p>
            <p>© {{ date.split()[-1] }} {{ company }}. All rights reserved.</p>
        </div>
    </div>
    
    <script src="https://cdn.jsdelivr.net/npm/bootstrap@5.1.3/dist/js/bootstrap.bundle.min.js"></script>
</body>
</html>
        '''
    
    def _create_matplotlib_chart(self, data: Any, metadata: Dict[str, Any]) -> str:
        """Create a chart using matplotlib and save as image."""
        plt.figure(figsize=(10, 6))
        
        chart_type = metadata.get('chart_type', 'line')
        
        if isinstance(data, pd.DataFrame):
            if chart_type == 'line':
                data.plot(kind='line')
            elif chart_type == 'bar':
                data.plot(kind='bar')
            elif chart_type == 'pie' and len(data.columns) == 1:
                data.plot(kind='pie', y=data.columns[0])
        
        plt.title(metadata.get('title', ''))
        plt.xlabel(metadata.get('xlabel', ''))
        plt.ylabel(metadata.get('ylabel', ''))
        
        # Save chart
        chart_path = os.path.join(self.config.output_path, f"chart_{datetime.now().timestamp()}.png")
        plt.savefig(chart_path, dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()
        
        return chart_path


class ReportFactory:
    """Factory for creating report generators."""
    
    @staticmethod
    def create_generator(config: ReportConfig) -> BaseReportGenerator:
        """
        Create appropriate report generator based on format.
        
        Args:
            config: Report configuration
            
        Returns:
            Report generator instance
        """
        generators = {
            ReportFormat.PDF: PDFReportGenerator,
            ReportFormat.EXCEL: ExcelReportGenerator,
            ReportFormat.POWERPOINT: PowerPointReportGenerator,
            ReportFormat.HTML: HTMLReportGenerator
        }
        
        generator_class = generators.get(config.format)
        if not generator_class:
            raise ValueError(f"Unsupported report format: {config.format}")
        
        return generator_class(config)


class FinancialReportBuilder:
    """High-level interface for building financial reports."""
    
    def __init__(self):
        """Initialize report builder."""
        self.data = {}
        self.analysis_results = {}
    
    def load_data(self, key: str, data: Any):
        """Load data for report."""
        self.data[key] = data
    
    def load_analysis(self, key: str, results: Any):
        """Load analysis results."""
        self.analysis_results[key] = results
    
    def build_portfolio_report(self, config: ReportConfig) -> str:
        """Build portfolio summary report."""
        generator = ReportFactory.create_generator(config)
        
        # Add title page
        generator.add_title_page()
        
        # Add executive summary
        summary_data = self._prepare_executive_summary()
        generator.add_executive_summary(summary_data)
        
        # Add portfolio overview
        if 'portfolio' in self.data:
            portfolio_df = self.data['portfolio']
            generator.add_section(ReportSection(
                title="Portfolio Holdings",
                content=portfolio_df,
                section_type="table"
            ))
        
        # Add performance metrics
        if 'performance' in self.analysis_results:
            perf_data = self.analysis_results['performance']
            generator.add_section(ReportSection(
                title="Performance Metrics",
                content=perf_data,
                section_type="table"
            ))
        
        # Add risk analysis
        if 'risk' in self.analysis_results:
            risk_data = self.analysis_results['risk']
            generator.add_section(ReportSection(
                title="Risk Analysis",
                content=risk_data,
                section_type="summary"
            ))
        
        # Add charts
        if config.include_charts:
            self._add_charts(generator)
        
        # Add disclaimer
        if config.include_disclaimer:
            generator.add_disclaimer()
        
        # Generate report
        return generator.generate()
    
    def build_risk_report(self, config: ReportConfig) -> str:
        """Build risk analysis report."""
        generator = ReportFactory.create_generator(config)
        
        # Add title page
        generator.add_title_page()
        
        # Add risk summary
        if 'risk_summary' in self.analysis_results:
            generator.add_section(ReportSection(
                title="Risk Summary",
                content=self.analysis_results['risk_summary'],
                section_type="summary"
            ))
        
        # Add VaR analysis
        if 'var_analysis' in self.analysis_results:
            generator.add_section(ReportSection(
                title="Value at Risk Analysis",
                content=self.analysis_results['var_analysis'],
                section_type="table"
            ))
        
        # Add stress test results
        if 'stress_tests' in self.analysis_results:
            generator.add_section(ReportSection(
                title="Stress Test Results",
                content=self.analysis_results['stress_tests'],
                section_type="table"
            ))
        
        # Add correlation analysis
        if 'correlations' in self.analysis_results:
            generator.add_section(ReportSection(
                title="Correlation Analysis",
                content=self.analysis_results['correlations'],
                section_type="chart",
                metadata={'chart_type': 'heatmap'}
            ))
        
        # Add disclaimer
        if config.include_disclaimer:
            generator.add_disclaimer()
        
        return generator.generate()
    
    def build_performance_report(self, config: ReportConfig) -> str:
        """Build performance analysis report."""
        generator = ReportFactory.create_generator(config)
        
        # Add title page
        generator.add_title_page()
        
        # Add performance summary
        if 'performance_summary' in self.analysis_results:
            generator.add_section(ReportSection(
                title="Performance Summary",
                content=self.analysis_results['performance_summary'],
                section_type="summary"
            ))
        
        # Add returns analysis
        if 'returns' in self.data:
            generator.add_section(ReportSection(
                title="Returns Analysis",
                content=self.data['returns'],
                section_type="chart",
                metadata={'chart_type': 'line', 'title': 'Cumulative Returns'}
            ))
        
        # Add attribution analysis
        if 'attribution' in self.analysis_results:
            generator.add_section(ReportSection(
                title="Performance Attribution",
                content=self.analysis_results['attribution'],
                section_type="table"
            ))
        
        # Add benchmark comparison
        if 'benchmark_comparison' in self.analysis_results:
            generator.add_section(ReportSection(
                title="Benchmark Comparison",
                content=self.analysis_results['benchmark_comparison'],
                section_type="table"
            ))
        
        # Add disclaimer
        if config.include_disclaimer:
            generator.add_disclaimer()
        
        return generator.generate()
    
    def _prepare_executive_summary(self) -> Dict[str, Any]:
        """Prepare executive summary data."""
        summary = {
            "Report Date": datetime.now().strftime("%B %d, %Y"),
            "Portfolio Value": "$1,234,567",
            "YTD Return": "12.5%",
            "Sharpe Ratio": "1.85",
            "Max Drawdown": "-8.2%",
            "Risk Level": "Moderate"
        }
        
        # Add custom summary data
        if 'summary' in self.analysis_results:
            summary.update(self.analysis_results['summary'])
        
        return summary
    
    def _add_charts(self, generator: BaseReportGenerator):
        """Add charts to report."""
        # Performance chart
        if 'returns' in self.data:
            generator.add_section(ReportSection(
                title="Performance Chart",
                content=self.data['returns'],
                section_type="chart",
                metadata={
                    'chart_type': 'line',
                    'title': 'Portfolio Performance',
                    'xlabel': 'Date',
                    'ylabel': 'Cumulative Return'
                }
            ))
        
        # Allocation chart
        if 'allocation' in self.data:
            generator.add_section(ReportSection(
                title="Portfolio Allocation",
                content=self.data['allocation'],
                section_type="chart",
                metadata={
                    'chart_type': 'pie',
                    'title': 'Asset Allocation'
                }
            ))
        
        # Risk chart
        if 'risk_metrics' in self.data:
            generator.add_section(ReportSection(
                title="Risk Metrics",
                content=self.data['risk_metrics'],
                section_type="chart",
                metadata={
                    'chart_type': 'bar',
                    'title': 'Risk Analysis',
                    'xlabel': 'Metric',
                    'ylabel': 'Value'
                }
            ))


def main():
    """Example usage of reports module."""
    # Create sample data
    dates = pd.date_range(start='2023-01-01', end='2023-12-31', freq='D')
    
    portfolio_data = pd.DataFrame({
        'Symbol': ['AAPL', 'GOOGL', 'MSFT', 'AMZN', 'BRK.B'],
        'Quantity': [1000, 500, 800, 300, 200],
        'Price': [150.25, 125.50, 380.75, 145.60, 350.25],
        'Value': [150250, 62750, 304600, 43680, 70050],
        'Weight (%)': [23.8, 9.9, 48.2, 6.9, 11.1]
    })
    
    performance_data = pd.DataFrame({
        'Metric': ['Total Return', 'Annualized Return', 'Volatility', 'Sharpe Ratio'],
        'Value': [0.125, 0.118, 0.185, 1.85]
    })
    
    # Create report builder
    builder = FinancialReportBuilder()
    builder.load_data('portfolio', portfolio_data)
    builder.load_analysis('performance', performance_data)
    
    # Generate PDF report
    pdf_config = ReportConfig(
        title="Portfolio Analysis Report",
        subtitle="Q4 2023 Performance Review",
        format=ReportFormat.PDF
    )
    
    pdf_path = builder.build_portfolio_report(pdf_config)
    print(f"PDF report generated: {pdf_path}")
    
    # Generate Excel report
    excel_config = ReportConfig(
        title="Portfolio Analysis Report",
        format=ReportFormat.EXCEL
    )
    
    excel_path = builder.build_portfolio_report(excel_config)
    print(f"Excel report generated: {excel_path}")
    
    # Generate HTML report
    html_config = ReportConfig(
        title="Portfolio Analysis Report",
        format=ReportFormat.HTML
    )
    
    html_path = builder.build_portfolio_report(html_config)
    print(f"HTML report generated: {html_path}")


if __name__ == "__main__":
    main()